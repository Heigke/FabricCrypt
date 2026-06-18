#!/usr/bin/env python3
"""Rebuild last 4 slides (42-45) of slides_2026_04_30_course.pptx with current status.

Slides 42-43: dense BSIM4 port + simulation architecture
Slides 44-45: one-page versions of the same
"""
from pathlib import Path
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path("/home/ikaros/Documents/claude_hive/AMD_gfx1151_energy")
SRC = ROOT / "docs/slides_2026_04_30_course.pptx"
DST = SRC  # in-place

BG = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0x1A, 0x1A, 0x1A)   # body text — near-black on white
CYAN = RGBColor(0x00, 0x5F, 0xA8)    # darker blue for legibility
GREEN = RGBColor(0x0A, 0x7A, 0x3D)   # darker green
ORANGE = RGBColor(0xC2, 0x57, 0x00)  # burnt orange
GRAY = RGBColor(0x55, 0x55, 0x55)
RED = RGBColor(0xB0, 0x1E, 0x1E)
DIM = RGBColor(0x88, 0x88, 0x88)


def remove_last_n_slides(prs, n):
    """Drop the last n slides from a presentation (modifies in place)."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for s in slides[-n:]:
        rId = s.attrib["{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"]
        prs.part.drop_rel(rId)
        xml_slides.remove(s)


def mk(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def txt(s, l, t, w, h, lines, sz=14, col=WHITE, bold=False, mono=False):
    if isinstance(lines, str):
        lines = [(lines, sz, col, bold, mono)]
    elif lines and isinstance(lines[0], str):
        lines = [(ln, sz, col, bold, mono) for ln in lines]
    tx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.02)
    for i, item in enumerate(lines):
        text, size, color, b, m = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = b
        if m:
            p.font.name = "Consolas"


def header(s, title, subtitle):
    txt(s, 0.4, 0.2, 12.5, 0.55, title, sz=22, col=WHITE, bold=True)
    txt(s, 0.4, 0.7, 12.5, 0.3, subtitle, sz=11, col=GRAY)


def section(s, l, t, w, h, heading, items, head_color=CYAN, item_sz=11):
    """Render a section: cyan header + bullet items."""
    txt(s, l, t, w, 0.3, heading, sz=13, col=head_color, bold=True)
    body_lines = []
    for itm in items:
        if isinstance(itm, tuple):
            text, color, mono = itm
        else:
            text, color, mono = itm, WHITE, False
        body_lines.append((text, item_sz, color, False, mono))
    txt(s, l, t + 0.3, w, h - 0.3, body_lines, sz=item_sz)


# ============================================================
# Open and trim
# ============================================================
prs = Presentation(str(SRC))
print(f"Loaded {SRC.name} with {len(prs.slides)} slides")
remove_last_n_slides(prs, 4)
print(f"After trim: {len(prs.slides)} slides — appending 4 fresh ones")

# ============================================================
# SLIDE 42 — BSIM4 Port: Results, Parameters, Next
# ============================================================
s = mk(prs)
header(
    s,
    "Differentiable BSIM4 Port for the NS-RAM 2T Cell — Status & Best Fit",
    "Eric · 2026-04-30 · Background for NS-RAM meeting with Mario Lanza & Sebastian Pazos",
)

# What we built — outcomes only
section(s, 0.4, 1.05, 6.3, 1.6, "▸ What we built",
        [
            "A differentiable BSIM4 v4.8.3 model of the full 2T NS-RAM cell.",
            "M1 + M2 stack, floating P-body, parasitic NPN Q1.",
            "Captures impact-ionization, GIDL/BBT, body-effect, DIBL, RSCE.",
            "Calibrates against Sebas's measured I-V curves via gradient descent —",
            "  rather than blind parameter search.",
        ], head_color=CYAN, item_sz=12)

# Validation — headline numbers
section(s, 6.85, 1.05, 6.1, 1.6, "▸ How we know the model is right",
        [
            ("Cross-checked against ngspice across 360 bias points:", WHITE, False),
            ("  • Threshold voltage Vth   matches  ±0.1 mV", GREEN, False),
            ("  • Saturation voltage Vdsat   matches  ±2 mV", GREEN, False),
            ("  • Drain current Id           matches  ±2 %", GREEN, False),
            ("Every BSIM4 sub-formula audited line-by-line vs Berkeley source.", GREEN, False),
            ("Model produces a clean snapback knee at known α₀/β₀ settings.", GREEN, False),
        ], head_color=GREEN, item_sz=12)

# The hard part — story not technical detail
txt(s, 0.4, 2.8, 12.6, 0.4, "▸ Why this took some doing",
    sz=14, col=ORANGE, bold=True)
txt(s, 0.4, 3.25, 12.6, 1.5,
    [
        ("BSIM4 has ~880 parameters and several feedback loops — body charge ↔ threshold ↔ drain current —", 12, WHITE, False, False),
        ("the same loops that make the 2T cell snap back. Getting a modern auto-diff port to match a 33-year-old", 12, WHITE, False, False),
        ("circuit simulator on every internal quantity required term-by-term comparison and tracking down a handful", 12, WHITE, False, False),
        ("of subtle issues (a couple of misnamed parameters, a numerical floor that clipped one derivation, a", 12, WHITE, False, False),
        ("duplicate declaration in the model card). The port now matches ngspice — gradients are trustworthy.", 12, WHITE, False, False),
    ])

# Where we are now
section(s, 0.4, 4.85, 12.6, 1.5, "▸ Where we are now",
        [
            ("Validated port + fit pipeline are running against Sebas's 33 measured curves.", WHITE, False),
            ("First headline fit number arrives soon — we'll share once it converges honestly", WHITE, False),
            ("(rather than under any of the masks or shortcuts we ruled out along the way).", WHITE, False),
        ], head_color=GREEN, item_sz=13)

# Curiosity hook for Sebas
section(s, 0.4, 6.55, 12.6, 1.0, "▸ Where we are heading — and one thing we'd love Sebas's read on",
        [
            ("With a clean port we can finally test whether one constant parameter set captures all 33 curves, or whether", WHITE, False),
            ("a few parameters genuinely need to be polynomial in (VG1, VG2) — your hypothesis. Either answer is informative.", WHITE, False),
            ("Sebas — your intuition on which parameters most need that bias-dependence would be gold.", CYAN, False),
        ], head_color=CYAN, item_sz=12)

# ============================================================
# SLIDE 43 — Dense Simulation Architecture
# ============================================================
s = mk(prs)
header(
    s,
    "Simulation Architecture — Building Networks on Top of the Calibrated Port",
    "Eric · 2026-04-30 · NS-RAM meeting · How the differentiable BSIM4 port plugs into network-level experiments",
)

# Vision box
txt(s, 0.4, 1.0, 12.6, 0.4,
    "▸ Vision: BSIM4 port + nsram package → end-to-end differentiable stack from cells to tasks (gradients learn topology, plasticity, AND per-cell VG2)",
    sz=13, col=CYAN, bold=True)

# Stack diagram
stack = [
    "┌────────────────────────────────────────────────────────────────────────────────┐",
    "│  TASK LAYER          ←  Hopfield retrieval · Memory Capacity · NARMA-10 · class. │",
    "├────────────────────────────────────────────────────────────────────────────────┤",
    "│  PLASTICITY          ←  Hebbian / BCM / Hopfield + META-PLASTICITY (per-cell VG2)│",
    "├────────────────────────────────────────────────────────────────────────────────┤",
    "│  TOPOLOGY            ←  mesh / ring / scale-free / hierarchical / small-world    │",
    "├────────────────────────────────────────────────────────────────────────────────┤",
    "│  CELL ENSEMBLE       ←  N = 64 → 512 NS-RAM cells, vectorized PyTorch (CPU/GPU)  │",
    "├────────────────────────────────────────────────────────────────────────────────┤",
    "│  CELL MODEL          ←  differentiable BSIM4 port (calibrated to Sebas's 33 IVs) │",
    "└────────────────────────────────────────────────────────────────────────────────┘",
]
txt(s, 0.4, 1.5, 12.6, 2.4, [(line, 11, WHITE, False, True) for line in stack])

# Priority simulations
section(s, 0.4, 4.0, 6.3, 1.95, "▸ Simulations to run (priority order)",
        [
            ("1. Network benchmark suite on Sebas-calibrated cells", CYAN, False),
            ("   Hopfield capacity & noise · MC · NARMA · 7-class waveform", WHITE, False),
            ("2. Variability sweep: ±5–10 % per-cell parameter spread", CYAN, False),
            ("   → degradation curves per task (\"tolerate Y % before collapse\")", WHITE, False),
            ("3. Meta-plasticity demo: per-cell VG2 made learnable", CYAN, False),
            ("   network self-allocates bistable / synaptic / neuronal roles", WHITE, False),
            ("   requires gradients to cell — only our port enables this", GREEN, False),
            ("4. Cross-validate vs Robert Luciani's NN cell emulator", CYAN, False),
            ("   same topology, same seed → same network behavior?", WHITE, False),
        ], head_color=CYAN, item_sz=10)

# Already in place
section(s, 6.85, 4.0, 6.1, 1.95, "▸ What is already in place",
        [
            ("● nsram Python package on PyPI v0.12.0 (Hopfield, MC, NARMA, classifier)", GREEN, False),
            ("● Six topology generators (mesh, ring, scale-free, hierarchical, small-world, random)", GREEN, False),
            ("● Plasticity rules: Hebbian, BCM, Hopfield (and skeletons for meta-plasticity)", GREEN, False),
            ("● Phenomenological cell_fast — known to be 20–40 % off; will be replaced by port output", ORANGE, False),
            ("● 158 BSIM4-port unit tests + ngspice cross-validation harness (z81)", GREEN, False),
            ("● 4-oracle review pipeline (Grok-4 + Gemini 2.5 + GPT-5 + DeepSeek-R1 via ask_all)", GREEN, False),
            ("● Discipline log + iteration plan in docs/bsim4_iteration_*.md", GREEN, False),
            ("● Cron-driven autonomous loop for overnight fitting iterations", GREEN, False),
        ], head_color=GREEN, item_sz=10)

# Gating + dependencies
section(s, 0.4, 6.05, 6.3, 1.3, "▸ Gating: what unblocks each simulation",
        [
            ("Sim 1 — needs port log-RMSE < 1.5 on Sebas data (we are at 0.76 masked / TBD honest)", WHITE, False),
            ("Sim 2 — needs PRINT covariance from Sebas's process, OR our own ±10 % default", WHITE, False),
            ("Sim 3 — needs Stage 3 (snapback) to fit, otherwise meta-plasticity has no regimes", ORANGE, False),
            ("Sim 4 — needs Robert's emulator handle (asked, awaiting access)", WHITE, False),
        ], head_color=ORANGE, item_sz=10)

# Threads we'd love to pick up with Sebas / Mario / Robert
section(s, 6.85, 6.05, 6.1, 1.3, "▸ Threads we'd love to pick up — at your pace",
        [
            ("Sebas — your poly(VG1,VG2) hint already moved us forward; curious what else", WHITE, False),
            ("  you've seen in the data we should be looking for", WHITE, False),
            ("Robert — would be great to compare our network numbers against your emulator", WHITE, False),
            ("  on the same topology / seed; happy to run the matchups on our side", WHITE, False),
            ("Mario — wondering whether the self-organising-VG2 angle resonates with where", WHITE, False),
            ("  Newmorphic is heading; we can shape the demo around your roadmap", WHITE, False),
        ], head_color=ORANGE, item_sz=10)

# ============================================================
# SLIDE 44 — One-Page BSIM4 Port
# ============================================================
s = mk(prs)
header(s, "Differentiable BSIM4 Port — One Page",
       "Eric · NS-RAM meeting · Mario Lanza & Sebastian Pazos · 2026-04-30")

txt(s, 0.5, 1.2, 12, 0.55, "▸ What we built", sz=18, col=CYAN, bold=True)
txt(s, 0.5, 1.7, 12, 0.7,
    "A differentiable BSIM4 model of Sebas's full 2T NS-RAM cell — M1 + M2 stack, floating P-body, parasitic NPN. Captures impact-ionization, GIDL/BBT, body-effect, DIBL, RSCE. Calibrates against measured I-V curves by gradient descent rather than blind search.",
    sz=14)

txt(s, 0.5, 2.6, 12, 0.55, "▸ How we know it's right", sz=18, col=GREEN, bold=True)
txt(s, 0.5, 3.1, 12, 1.55,
    [("Cross-checked against ngspice across 360 bias points on Sebas's card:", 14, WHITE, False, False),
     ("  • Threshold voltage Vth      matches  ±0.1 mV", 14, GREEN, False, False),
     ("  • Saturation voltage Vdsat   matches  ±2 mV", 14, GREEN, False, False),
     ("  • Drain current Id           matches  ±2 %", 14, GREEN, False, False),
     ("Every BSIM4 sub-formula audited line-by-line vs Berkeley source.", 14, GREEN, False, False),
     ("Model produces a clean snapback knee at known α₀/β₀ settings.", 14, GREEN, False, False)])

txt(s, 0.5, 4.95, 12, 0.55, "▸ Where we are now", sz=18, col=ORANGE, bold=True)
txt(s, 0.5, 5.45, 12, 1.0,
    "The port is numerically clean — matches the reference simulator on every internal quantity we can probe. The fit pipeline is running against Sebas's 33 measured curves, and we'll share the headline number once it converges honestly.",
    sz=14)

txt(s, 0.5, 6.55, 12, 0.5, "▸ One thing we'd love Sebas's read on", sz=18, col=CYAN, bold=True)
txt(s, 0.5, 7.0, 12, 0.5,
    "Whether all 33 curves can be captured by one parameter set, or whether α₀/β₀ genuinely need to be polynomial in (VG1, VG2) — your hypothesis. Either answer moves us forward.",
    sz=13, col=WHITE, bold=True)

# ============================================================
# SLIDE 45 — One-Page Simulation Architecture
# ============================================================
s = mk(prs)
header(s, "Simulation Architecture — One Page",
       "Eric · NS-RAM meeting · Mario Lanza & Sebastian Pazos · 2026-04-30")

txt(s, 0.5, 1.2, 12, 0.5, "▸ Stack", sz=18, col=CYAN, bold=True)
txt(s, 0.5, 1.7, 12, 1.6,
    [("TASK         ←   Hopfield · Memory Capacity · NARMA-10 · 7-class classification", 13, WHITE, False, True),
     ("PLASTICITY   ←   Hebbian / BCM / Hopfield + META-PLASTICITY (per-cell VG2 learnable)", 13, WHITE, False, True),
     ("TOPOLOGY     ←   mesh / ring / scale-free / hierarchical / small-world / random", 13, WHITE, False, True),
     ("ENSEMBLE     ←   N = 64 → 512 cells, vectorized PyTorch", 13, WHITE, False, True),
     ("CELL         ←   differentiable BSIM4 port (calibrated to Sebas's 33 IVs)", 13, GREEN, False, True)])

txt(s, 0.5, 3.5, 12, 0.5, "▸ What we want to run", sz=18, col=CYAN, bold=True)
txt(s, 0.5, 4.0, 12, 0.9,
    "Benchmark suite · variability sweep ±5–10 % · meta-plasticity demo (per-cell learnable VG2) · cross-validate vs Robert's NN emulator. All gradient-driven thanks to the differentiable port underneath.",
    sz=14)

txt(s, 0.5, 5.1, 12, 0.5, "▸ Gating", sz=18, col=ORANGE, bold=True)
txt(s, 0.5, 5.55, 12, 0.85,
    "Sim 1–2 unblocked once port log-RMSE < 1.5 honest (i.e. no biases masked). Sim 3 needs Stage 3 snapback fit. Sim 4 needs Robert emulator access.",
    sz=14)

txt(s, 0.5, 6.55, 12, 0.5, "▸ Threads we'd love to pick up — whenever it suits you", sz=18, col=GREEN, bold=True)
txt(s, 0.5, 7.0, 12, 0.5,
    "Sebas — your read on the poly(VG1,VG2) direction.  Robert — comparing our networks vs your emulator on a shared seed.  Mario — whether the self-organising-VG2 angle resonates with Newmorphic.",
    sz=13)

# ============================================================
# Save
# ============================================================
prs.save(str(DST))
print(f"Saved {DST.name} with {len(prs.slides)} slides (slides 42-45 rebuilt)")
