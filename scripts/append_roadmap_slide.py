#!/usr/bin/env python3
"""Append Eric's roadmap slide to the combined Mario/Seb deck.

Inserts as slide 5 (right after Eric's existing slides 2-4 and before
Robert's image slides).
"""
from pathlib import Path
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

DECK = Path("/home/ikaros/Documents/claude_hive/AMD_gfx1151_energy/docs/NSRAM 20260429 Mario Seb.pptx")

# Match the cleaned, white-background style from slide 2/4
BG = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
CYAN = RGBColor(0x00, 0x5F, 0xA8)
GREEN = RGBColor(0x0A, 0x7A, 0x3D)
ORANGE = RGBColor(0xC2, 0x57, 0x00)
GRAY = RGBColor(0x55, 0x55, 0x55)
DIM = RGBColor(0x88, 0x88, 0x88)


def mk_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def txt(s, l, t, w, h, lines, sz=12, col=TEXT, bold=False, mono=False):
    if isinstance(lines, str):
        lines = [(lines, sz, col, bold, mono)]
    elif lines and isinstance(lines[0], str):
        lines = [(ln, sz, col, bold, mono) for ln in lines]
    tx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.02)
    for i, item in enumerate(lines):
        text, size, color, b, m = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = b
        if m:
            p.font.name = "Consolas"


def header(s, title, subtitle):
    txt(s, 0.4, 0.2, 12.5, 0.55, title, sz=22, col=TEXT, bold=True)
    txt(s, 0.4, 0.7, 12.5, 0.3, subtitle, sz=11, col=GRAY)


def section(s, l, t, w, h, heading, items, head_color=CYAN, item_sz=11):
    txt(s, l, t, w, 0.3, heading, sz=13, col=head_color, bold=True)
    body = []
    for itm in items:
        if isinstance(itm, tuple):
            text, color, mono = itm
        else:
            text, color, mono = itm, TEXT, False
        body.append((text, item_sz, color, False, mono))
    txt(s, l, t + 0.3, w, h - 0.3, body, sz=item_sz)


def move_slide_to_position(prs, slide_idx_from, position_to):
    """Move a slide from current index to target position (0-indexed)."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    moved = slides[slide_idx_from]
    xml_slides.remove(moved)
    xml_slides.insert(position_to, moved)


def main():
    prs = Presentation(str(DECK))
    n_before = len(prs.slides)
    print(f"Loaded deck with {n_before} slides")

    # Build the roadmap slide (initially at the end)
    s = mk_slide(prs)
    header(
        s,
        "Network/Architecture Roadmap — Eric, next 9 months",
        "Eric · 2026-04-30 · Complementary to Robert's emulator track · Subject to your priorities",
    )

    # Top row — Q2 2026 + Q3 2026
    section(s, 0.4, 1.05, 6.3, 1.95, "▸ Q2 2026 (May–June) — Calibration locked",
            [
                ("Finalize fit pipeline on cleaned card; lock arclength solver as default.", TEXT, False),
                ("First honest log-RMSE on Sebas's 33 curves (constant params).", TEXT, False),
                ("If single param set insufficient → polynomial(VG1, VG2) fit per your hypothesis.", TEXT, False),
                ("Cross-validate fit numerically against Robert's Julia emulator.", TEXT, False),
                ("Deliverable: calibrated nsram package v0.13 with reproducible fit.", GREEN, False),
            ], head_color=CYAN, item_sz=11)

    section(s, 6.85, 1.05, 6.1, 1.95, "▸ Q3 2026 (July–Sept) — Network benchmarks",
            [
                ("Hopfield retrieval + Memory Capacity + NARMA-10 + 7-class waveform.", TEXT, False),
                ("N = 64 → 512 cells scaling study on 6 topologies.", TEXT, False),
                ("Variability sweep ±5–10 % per-cell → degradation curves per task.", TEXT, False),
                ("Output: \"design rule\" doc — for task X, you can tolerate Y % spread.", TEXT, False),
                ("Deliverable: standardized network benchmark suite (Sebas + Robert).", GREEN, False),
            ], head_color=CYAN, item_sz=11)

    # Middle row — Q4 2026 + Q1 2027
    section(s, 0.4, 3.15, 6.3, 1.95, "▸ Q4 2026 (Oct–Dec) — Meta-plasticity demo",
            [
                ("Per-cell VG2 made learnable; multi-task signal drives self-allocation.", TEXT, False),
                ("Network discovers its own bistable / synaptic / neuronal split.", TEXT, False),
                ("Requires gradients all the way to the cell — only the differentiable port enables it.", TEXT, False),
                ("Headline claim: \"no memristor can do this\" (uniquely NS-RAM).", ORANGE, False),
                ("Deliverable: paper-ready demo + ablation study.", GREEN, False),
            ], head_color=CYAN, item_sz=11)

    section(s, 6.85, 3.15, 6.1, 1.95, "▸ Q1 2027 (Jan–Mar) — Joint paper + tape-out input",
            [
                ("Draft: \"Array-scale computation in CMOS NS-RAM: from single-cell", TEXT, False),
                ("  physics to N-cell reservoir with self-organized VG2\".", TEXT, False),
                ("Cell-parameter ranking by task sensitivity → input to next tape-out.", TEXT, False),
                ("Recommended VG2 sweep range + variability tolerance per architecture.", TEXT, False),
                ("Deliverable: paper draft + technical note for Sebas/Mario.", GREEN, False),
            ], head_color=CYAN, item_sz=11)

    # Risks + asks
    section(s, 0.4, 5.25, 6.3, 1.55, "▸ Known risks + buffers",
            [
                ("Port has 10–15 % saturation residual on short-channel M1 (skipped sub-blocks).", ORANGE, False),
                ("  → buffer 1–2 weeks in Q2 to port velocity-overshoot if it blocks fit.", DIM, False),
                ("Polynomial param form ~1–2 weeks to implement if Q2 needs it.", ORANGE, False),
                ("Network sims gated on Robert's emulator handle — Q2 dependency.", ORANGE, False),
            ], head_color=ORANGE, item_sz=11)

    section(s, 6.85, 5.25, 6.1, 1.55, "▸ Things we'd love your input on (no rush)",
            [
                ("Sebas — poly(VG1, VG2) form & degree; transient Id(t) for τ-calibration.", TEXT, False),
                ("Robert — emulator handle for cross-validation; weekly cadence?", TEXT, False),
                ("Mario — does meta-plasticity narrative resonate with Newmorphic strategy?", TEXT, False),
                ("Mario — array-size target for next tape-out shapes our Q3 scaling study.", TEXT, False),
            ], head_color=GREEN, item_sz=11)

    # Bottom — closing line
    section(s, 0.4, 6.95, 12.6, 0.7, "▸ Spirit of the roadmap",
            [
                ("Each quarter has one tangible deliverable that's useful to Sebas's lab and to Mario's roadmap regardless", TEXT, False),
                ("of whether the next quarter lands. Re-prioritize freely — this is a proposal, not a commitment.", CYAN, False),
            ], head_color=CYAN, item_sz=11)

    # Move new slide from end (index n_before) to position 4 (= slide 5, after slide 4)
    move_slide_to_position(prs, slide_idx_from=n_before, position_to=4)

    out_path = DECK
    prs.save(str(out_path))
    print(f"Saved {out_path}  (now {len(prs.slides)} slides; roadmap is slide 5)")


if __name__ == "__main__":
    main()
