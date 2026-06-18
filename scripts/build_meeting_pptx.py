"""Build the 2026-04-30 NS-RAM meeting slide deck as a real .pptx file
with embedded result images.

Run: venv/bin/python -m scripts.build_meeting_pptx
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = Path("/home/ikaros/Documents/claude_hive/AMD_gfx1151_energy")
OUT = ROOT / "docs" / "slides_2026_04_30_course.pptx"

NAVY    = RGBColor(0x1F, 0x3A, 0x5F)
ORANGE  = RGBColor(0xE6, 0x7E, 0x22)
GRAY    = RGBColor(0x55, 0x55, 0x55)
LIGHT   = RGBColor(0xEC, 0xF0, 0xF1)
RED     = RGBColor(0xC0, 0x39, 0x2B)
GREEN   = RGBColor(0x27, 0xAE, 0x60)


def add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Background fill (left navy strip)
    from pptx.oxml.ns import qn
    from copy import deepcopy
    # Title
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(2.0),
                                   Inches(12.0), Inches(2.0))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = NAVY
    # Subtitle
    sx = slide.shapes.add_textbox(Inches(0.6), Inches(4.2),
                                   Inches(12.0), Inches(1.5))
    sf = sx.text_frame
    p = sf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24); p.font.color.rgb = GRAY; p.font.italic = True
    return slide


def add_section_slide(prs, section_no, section_title):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    # Big number
    nx = slide.shapes.add_textbox(Inches(0.6), Inches(2.0),
                                   Inches(4.0), Inches(3.0))
    nf = nx.text_frame
    p = nf.paragraphs[0]
    p.text = section_no
    p.font.size = Pt(120); p.font.bold = True; p.font.color.rgb = ORANGE
    # Title
    tx = slide.shapes.add_textbox(Inches(4.5), Inches(3.0),
                                   Inches(8.5), Inches(2.0))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = NAVY
    return slide


def add_bullets_slide(prs, title, bullets, *, image=None, image_caption=None):
    """bullets: list of (text, level) tuples or just strings (level=0)."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    # Title bar
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.25),
                                   Inches(12.5), Inches(0.85))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = NAVY
    # underline bar
    from pptx.shapes.connector import Connector
    line = slide.shapes.add_connector(1, Inches(0.4), Inches(1.05),
                                       Inches(12.9), Inches(1.05))
    line.line.color.rgb = ORANGE; line.line.width = Pt(2)
    # Body
    if image is not None:
        body_w = Inches(7.0); body_h = Inches(5.5)
    else:
        body_w = Inches(12.5); body_h = Inches(5.8)
    bx = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), body_w, body_h)
    bf = bx.text_frame; bf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        if i == 0:
            p = bf.paragraphs[0]
        else:
            p = bf.add_paragraph()
        p.level = level
        # bullet character
        prefix = "▸ " if level == 0 else "• "
        p.text = prefix + text
        p.font.size = Pt(20 - level * 3)
        if level == 0:
            p.font.color.rgb = NAVY; p.font.bold = False
        else:
            p.font.color.rgb = GRAY
        p.space_after = Pt(8)
    # Image
    if image is not None and image.exists():
        ix = slide.shapes.add_picture(str(image), Inches(7.7), Inches(1.4),
                                       width=Inches(5.5))
        if image_caption:
            cx = slide.shapes.add_textbox(Inches(7.7), Inches(6.7),
                                           Inches(5.5), Inches(0.5))
            cf = cx.text_frame
            p = cf.paragraphs[0]
            p.text = image_caption
            p.font.size = Pt(11); p.font.color.rgb = GRAY; p.font.italic = True
            p.alignment = PP_ALIGN.CENTER
    return slide


def add_diagram_slide(prs, title, mono_text, footer=None):
    """Slide with a monospace ASCII diagram."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.25),
                                   Inches(12.5), Inches(0.85))
    tf = tx.text_frame
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = NAVY
    line = slide.shapes.add_connector(1, Inches(0.4), Inches(1.05),
                                       Inches(12.9), Inches(1.05))
    line.line.color.rgb = ORANGE; line.line.width = Pt(2)
    bx = slide.shapes.add_textbox(Inches(0.5), Inches(1.3),
                                   Inches(12.3), Inches(5.5))
    bf = bx.text_frame; bf.word_wrap = False
    lines = mono_text.split("\n")
    for i, ln in enumerate(lines):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = ln
        p.font.name = "Consolas"
        p.font.size = Pt(16)
        p.font.color.rgb = NAVY
    if footer:
        fx = slide.shapes.add_textbox(Inches(0.5), Inches(6.7),
                                       Inches(12.3), Inches(0.5))
        ff = fx.text_frame
        p = ff.paragraphs[0]
        p.text = footer
        p.font.size = Pt(14); p.font.color.rgb = GRAY; p.font.italic = True
    return slide


# ─────────────────────────────────────────────────────────────────────
# Build the deck
# ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── 1. Title ──
    add_title_slide(prs,
        "NS-RAM From Scratch",
        "A Course for Anyone Stepping Into This Project   ·   2026-04-30   ·   Eric & Claude")

    # ── 2. How to read these slides ──
    add_bullets_slide(prs, "How to Use These Slides", [
        "Each slide stands alone. Read top-to-bottom.",
        ("By the end you'll understand:", 0),
        ("What a transistor is and why it limits modern computing", 1),
        ("What 'memory devices' research is exploring", 1),
        ("What NS-RAM is and why Mario thinks it's special", 1),
        ("Who in our group is doing what", 1),
        ("Where we (Eric's side) fit, and what we're proposing", 1),
        ("Time to read: ~25 min   ·   Time to skim: ~10 min", 0),
    ])

    # ── PART I ──
    add_section_slide(prs, "I", "From Transistors to Memory Cells")

    add_diagram_slide(prs, "Slide 1 — What is a Transistor?",
"""    Source ──────┤├────── Drain
                  │
                [Gate]

    Source / Drain : where current flows through
    Gate           : control wire — voltage decides if current flows
""",
        footer="Analogy: a kitchen tap. Pipe ends = source/drain. Hand on tap = gate. "
               "A modern chip has ~50 billion of these on a fingernail.")

    add_diagram_slide(prs, "Slide 2 — The MOSFET (the workhorse)",
"""                       Gate
                        |
       Oxide layer  ────┼────    (insulator)
                        v
       ─Source─[ channel ]─Drain─    (silicon)
                        |
                       Body            (silicon underneath)

    Gate voltage  →  electric field through oxide
                  →  pulls electrons into a thin channel
                  →  current can flow source → drain
""",
        footer="Analogy: magnet over sand. Iron filings line up. "
               "Threshold voltage Vth = field needed before channel exists.")

    add_bullets_slide(prs, "Slide 3 — How We Build Memory From Transistors", [
        "DRAM — 1 transistor + 1 capacitor",
        ("Capacitor = leaky bucket holding a few electrons", 1),
        ("Bit = 'is the bucket full?'", 1),
        ("Leaks every ~64 ms → must refresh constantly", 1),
        ("Cheap and dense, but volatile", 1),
        "SRAM — 6 transistors, cross-coupled inverters",
        ("Two inverters fighting each other → stays in one of two states", 1),
        ("Fast, stable while powered", 1),
        ("Big, expensive (6× the area per bit)", 1),
        "Both lose ALL data when power is off.",
    ])

    add_bullets_slide(prs, "Slide 4 — Why People Want Something Better", [
        "Three pain points:",
        ("Memory wall: moving data costs more energy than computing", 1),
        ("Volatility: every reboot from scratch", 1),
        ("Brain envy: your brain runs on 20 W, learns continuously, never reboots", 1),
        "Idea: what if memory cells could also COMPUTE?",
        ("And remember without power?", 1),
        ("And learn?", 1),
        "Field of emerging non-volatile memory: RRAM, PCM, MRAM, memristors, NS-RAM",
    ])

    add_diagram_slide(prs, "Slide 5 — The Crossbar Architecture",
"""               col 1   col 2   col 3
                |       |       |
     row 1 ─────●───────●───────●─────
                |       |       |
     row 2 ─────●───────●───────●─────
                |       |       |
     row 3 ─────●───────●───────●─────

    Each ● stores a number W[i,j] (its conductance).
    Send voltages on rows  →  currents on columns = V × W.

    That's matrix multiplication, done in physics, no transport.
""",
        footer="AI is mostly matrix multiplication. Imagine doing it in memory itself.")

    # ── PART II ──
    add_section_slide(prs, "II", "The Cast of Memory Devices")

    add_bullets_slide(prs, "Slide 6 — Memristors: The Famous Star", [
        "A memristor changes its conductance based on voltage history. "
        "Layer of metal-oxide where ions move around.",
        "Pros:",
        ("Tiny, dense, non-volatile", 1),
        ("Can be made analog (continuous conductance)", 1),
        "Cons (the elephant in the room):",
        ("Variability: ±20–50% spread per cell", 1),
        ("Yield: many cells just don't work after fabrication", 1),
        ("Drift: conductance changes over time", 1),
        ("Endurance: dies after ~10⁶ writes", 1),
        "Analogy: vinyl records pressed by a slightly drunk machinist. "
        "You can't ship a CPU made of those.",
    ])

    add_bullets_slide(prs, "Slide 7 — Brain Envy: Two Modes a Cell Could Have", [
        "Your brain has two kinds of components:",
        "Neurons — they SPIKE. Information lives in timing.",
        ("'I just fired'  — like a binary alarm clock that resets", 1),
        "Synapses — they STORE STRENGTH. Information lives in how much they pass.",
        ("'When neuron A fires, I let through 0.7× of the signal'", 1),
        ("Like a continuously adjustable volume knob", 1),
        "A traditional chip has ONE element (the transistor).",
        ("A neural-inspired chip ideally has BOTH.", 1),
        "What if ONE device could do both, and switch between them?",
    ])

    # ── PART III ──
    add_section_slide(prs, "III", "NS-RAM: The New Kid")

    add_diagram_slide(prs, "Slide 8 — What is NS-RAM?",
"""        VG1 (top gate)
          |
        ┌─┴─┐
    ────┤M1 ├────
        └───┘
        [BODY]    ← floating, holds charge, no direct wire
        ┌───┐
    ────┤M2 ├────
        └─┬─┘
          |
         VG2 (bottom gate)

    Two stacked MOSFETs sharing a 'floating' silicon body in between.
    Standard CMOS process — same fab line as your phone's chip.
""",
        footer="Developed at NUS by Mario Lanza & Sebastian Pazos. "
               "Two control gates VG1 and VG2; the body between them is isolated.")

    add_bullets_slide(prs, "Slide 9 — Why the Floating Body Matters", [
        "Body holds a few thousand electrons. Like a tiny battery: "
        "charges in 10 ms, leaks in ~10 seconds.",
        "Two stable states:",
        ("Body charge HIGH  →  threshold low   →  channel ON   →  reads '1'", 1),
        ("Body charge LOW   →  threshold high  →  channel OFF  →  reads '0'", 1),
        "Bistability — body sits in one of two energy wells.",
        ("Flipping requires an event (an 'avalanche' inside the silicon)", 1),
        "Analogy: a ball in a landscape with two valleys. "
        "Most of the time it stays in one valley. Hit it hard → jumps to the other.",
    ])

    add_diagram_slide(prs, "Slide 10 — Snapback: The Switching Mechanism",
"""   Id |
      |           ┌─── snap! current jumps
      |          /
      |         /
      |  ──────         (gradual rise)
      |
      └────────────────  Vd

    Increase drain voltage slowly → current rises gently
    Then SUDDENLY: an electron-hole avalanche dumps charge into the body
                   body switches state → cell LATCHES.

    To unlatch: pull the body charge out via VG2.
""",
        footer="That's how you write a '1' or '0'. It can also be made gradual — that gives the synaptic mode.")

    add_diagram_slide(prs, "Slide 11 — The Killer Feature: VG2 Selects the Regime",
"""    VG2 = -0.2 V          VG2 ≈ 0 V             VG2 = +0.4 V

      ┌────────┐           ┌────────┐            ┌────────┐
      |  ┌─┐   |           |  ╱╲    |            |    ╱   |
      |  | |   |           | ╱  ╲   |            |   ╱    |
      | ─┘ └── |           | ─    ─ |            |  ─     |
      └────────┘           └────────┘            └────────┘
       BISTABLE             SOFT/SYNAPSE          INTEGRATOR
       (memory)             (analog weight)       (neuron-like)
""",
        footer="Same physical cell. Same silicon. The bottom gate (VG2) chooses the mode. No memristor can do this.")

    add_bullets_slide(prs, "Slide 12 — Why This Matters Industrially", [
        "Three things distinguish NS-RAM from memristors:",
        "CMOS-compatible — built in same fab as standard chips, no exotic materials",
        "100% yield + low variability (Mario's Nature 2025)",
        ("Every cell works. ±5–10% spread vs ±50% for memristors", 1),
        "Dual-mode via VG2 — one device, three behaviors, switchable in real time",
        "Translation: this could actually scale to billions of cells in a real foundry.",
        ("Memristors couldn't.", 1),
    ])

    # ── PART IV ──
    add_section_slide(prs, "IV", "Who Does What")

    add_bullets_slide(prs, "Slide 13 — Mario Lanza: The Principal Investigator", [
        "Role: PI of the NS-RAM program at NUS. Vision-setter, fundraiser, paper author.",
        "What he owns:",
        ("The big strategic narrative ('NS-RAM is the post-memristor')", 1),
        ("Publications (Nature 2025 was his)", 1),
        ("Funding, foundry relationships, industrial partnerships", 1),
        "What he wants from collaborators:",
        ("Results that reinforce the narrative", 1),
        ("Practical demonstrations he can show investors", 1),
        ("Designs that could move toward a chip", 1),
        "Analogy: the orchestra conductor. Doesn't play one instrument, makes the whole thing sound coherent.",
    ])

    add_bullets_slide(prs, "Slide 14 — Sebastian Pazos: The Device Physicist", [
        "Role: the cell guy. Fabricates them, measures them, fits SPICE models.",
        "What he owns:",
        ("Physical wafers from foundry (130 nm node currently)", 1),
        ("I-V measurement data (we just got CSVs)", 1),
        ("BSIM4 + parasitic-NPN SPICE deck (foundry-calibrated)", 1),
        ("New model fit (in progress, swapped to BSIM4 impact-ionization)", 1),
        "Recent update from him:",
        ("Dropped the avalanche-diode model (numerically painful)", 1),
        ("Now using BSIM4 §6.1 (impact-ionization) + complementary bipolar current", 1),
        ("Sent us 2T schematic + raw I-V curves + 130nm PTM card", 1),
        "Analogy: the chemist. Knows EXACTLY what's in the test tube.",
    ])

    add_bullets_slide(prs, "Slide 15 — Robert Luciani: The Numerical Surrogate Guy", [
        "Role: building a fast, differentiable simulator that mimics SPICE.",
        "What he owns:",
        ("A neural-net emulator trained on SPICE outputs", 1),
        ("~30 000× faster than SPICE", 1),
        ("End-to-end differentiable (you can backprop through it)", 1),
        ("Now also porting things to a Julia simulator", 1),
        "Why it matters: SPICE is slow.",
        ("1000 cells × 10 s × 1 µs resolution: SPICE = hours, Robert's = seconds", 1),
        "Analogy: Robert built a video codec for SPICE simulations. Looks the same, plays 30 000× faster.",
    ])

    add_bullets_slide(prs, "Slide 16 — Eric & Claude (Us): The Network/Architecture Layer", [
        "Role: what happens when you wire many NS-RAM cells together?",
        "What we own:",
        ("nsram Python package on PyPI (v0.12.0)", 1),
        ("Phenomenological cell model (cell_fast, ~10 000× faster than SPICE)", 1),
        ("Six topology generators (ring, scale-free, hierarchical, ...)", 1),
        ("Plasticity rules: Hebbian, BCM, Hopfield", 1),
        ("An honest record of what worked and what didn't (z31–z57)", 1),
        "The honest gap:",
        ("We aren't competing with Sebas on cell fitting (we tried, z31–z37 stalled)", 1),
        ("We aren't competing with Robert on cell emulation (his is more accurate)", 1),
        "We are one level up: 'what computes when N cells are connected?'",
    ])

    add_diagram_slide(prs, "Slide 17 — One Sentence Each",
"""    Sebastian:  'I make and measure the cell.'

    Robert:     'I make a fast, differentiable copy of the cell.'

    Mario:      'I tell the world why the cell matters and get it funded.'

    Us (Eric):  'I wire many cells together and study what they compute.'

    ┌────────────────────────────────┐
    │  NETWORK / APPLICATION         │  ← Eric
    ├────────────────────────────────┤
    │  FAST CELL EMULATOR            │  ← Robert
    ├────────────────────────────────┤
    │  PHYSICAL CELL + SPICE         │  ← Sebastian
    └────────────────────────────────┘
       (Mario sets strategy around it all)
""",
        footer="Complementary, not competing. They form a stack.")

    # ── PART V ──
    add_section_slide(prs, "V", "What We Tried, What We Learned")

    add_bullets_slide(prs, "Slide 18 — Network-Level Findings That Survived Robustness Testing", [
        "Hopfield associative memory works on bistable cells",
        ("N=96, K=8 patterns, 25% input noise → 100% recall", 1),
        ("At K=25 (above 0.14·N textbook limit): bistable ~89% vs linear ~82%", 1),
        "Sparse connectivity beats fully-connected",
        ("Random / ring / small-world / scale-free / hierarchical: MC ≈ 3.7–4.1", 1),
        ("Fully-connected: MC drops to 2.8", 1),
        "Optimal regime is SUBCRITICAL",
        ("Branching ratio m ≈ 0.57 at peak memory (not edge-of-chaos!)", 1),
    ], image=ROOT / "results" / "z51_hopfield_memory" / "hopfield.png",
       image_caption="z51: Hopfield capacity, noise tolerance, bistable vs linear")

    add_bullets_slide(prs, "Slide 19 — Hopfield Capacity: Bistable vs Linear", [
        "Sweep number of stored patterns K from 5 to 50 at 25% input noise.",
        "Below K=15: both substrates achieve perfect recall.",
        "Above K=40: both saturate.",
        "K=20–35: bistable cells show graceful-degradation advantage.",
        ("K=25: bistable 89% vs linear 82% (+7 pp)", 1),
        ("K=35: bistable 73% vs linear 65% (+8 pp)", 1),
        "Bistability is NOT extra capacity — it's INSURANCE in the overloaded regime.",
    ], image=ROOT / "results" / "z52_hopfield_capacity_diff" / "capacity_diff.png",
       image_caption="z52: capacity difference between bistable and linear cells")

    add_bullets_slide(prs, "Slide 20 — Robustness to Device Variability", [
        "Foundry parts will have ~5–10% device-to-device spread.",
        "Test: per-cell Gaussian variability on K_back, A_iii, V_bjt_on.",
        "Bistable advantage SURVIVES realistic variability:",
        ("K=25: ~3–4 pp gap stable from 0% to 20% spread", 1),
        ("Both substrates degrade gracefully — no cliff", 1),
        "This is the key strength NS-RAM offers over memristors.",
    ], image=ROOT / "results" / "z53_hopfield_variability" / "variability.png",
       image_caption="z53: recall accuracy vs device parameter spread")

    add_bullets_slide(prs, "Slide 21 — Bistability Helps Most in Small Arrays", [
        "Scale N from 64 to 384 at fixed K/N = 0.25.",
        "N=64: bistable advantage = +7.5 pp",
        "N=96–128: +2.5 pp",
        "N≥192: both saturate to ~95–100%, advantage vanishes",
        "Concrete recommendation:",
        ("Bistable NS-RAM helps most for SMALL arrays (under ~128 cells)", 1),
        ("For large arrays (1000+) a linear substrate is essentially as capable", 1),
        ("for Hopfield-style retrieval", 1),
    ], image=ROOT / "results" / "z54_hopfield_scaling" / "scaling.png",
       image_caption="z54: bistable-vs-linear gap shrinks with array size")

    add_bullets_slide(prs, "Slide 22 — What We Retracted (and Why)", [
        "We had to walk back THREE earlier claims after stricter tests:",
        "'Heterogeneous α gives +72% MC' → was vs a weak baseline.",
        ("With strong baseline, heterogeneity LOSES by 11%", 1),
        "'Small-world is uniquely best' → after calibration, random barely edges it out.",
        ("Signal is sparse-vs-dense, not topology flavor", 1),
        "'VG2 sweet-spot at +0.20 V' → after K_back calibration, optimum flipped to −0.20 V",
        "Why I'm telling you: science benefits from this kind of correction.",
        ("Better to walk back at this stage than at peer review", 1),
    ])

    add_bullets_slide(prs, "Slide 23 — The Big Honest Caveat", [
        "Our cell model (cell_fast) is a phenomenological approximation.",
        ("NOT foundry-calibrated", 1),
        "We tried full BSIM4 fitting (z31–z37) → stalled at ~1.65 decade residual",
        "Pivoted to a 4-equation model — fast but off by ~20–40% on snapback",
        "VG2 polarity in our model DISAGREES with Sebas's measurements",
        ("We wrote a question to clarify (in docs/sebas_vg2_question.md)", 1),
        "Implication:",
        ("Network-level RANKINGS should be qualitatively right", 1),
        ("MAGNITUDES will shift on Robert's emulator or Sebas's SPICE", 1),
        "We have to say this openly. Otherwise our results look stronger than they are.",
    ])

    add_bullets_slide(prs, "Slide 24 — Yesterday's Result: Cell-Property Sensitivity", [
        "Swept 4 cell-level knobs, measured impact on 3 tasks.",
        "Robust finding: snapback strength is a DESIGN TENSION:",
        ("Low A_iii → good for CAM (Hopfield)", 1),
        ("High A_iii → good for reservoir (Memory Capacity)", 1),
        "Surprise that overturns z52:",
        ("With T_relax=300 (vs 50), LINEAR baseline beats bistable on Hopfield by 35%", 1),
        ("Bistable advantage has a regime caveat we missed before", 1),
        "Net effect: more honesty, more nuance, better story.",
    ], image=ROOT / "results" / "z57_cell_ablation_v2" / "ablation_v2.png",
       image_caption="z57: knob sensitivity heatmap + per-task response curves")

    # ── PART VI ──
    add_section_slide(prs, "VI", "The New Idea: Self-Organizing Plasticity")

    add_diagram_slide(prs, "Slide 25 — The Big Idea",
"""    Recall: VG2 selects the cell's regime (bistable / synapse / neuron).

    What if the network DECIDED which cells should be in which mode?

    Before training:                  After meta-plasticity:

      ┌────────────────────────┐       ┌────────────────────────┐
      | all cells: VG2 = -0.10 |       | ~30%  bistable (-0.20) |
      |        (uniform)       |       | ~50%  synaptic ( 0.00) |
      |                        |       | ~20%  neuronal (+0.40) |
      └────────────────────────┘       └────────────────────────┘
""",
        footer="The network learns its own division of labor. No memristor can do this. "
               "Uniquely-NS-RAM thing.")

    add_bullets_slide(prs, "Slide 26 — Why This Is The Right Pitch", [
        "Three reasons it's strong for the meeting:",
        "Uses NS-RAM's actual unique feature (dual-mode VG2)",
        ("instead of competing with memristor-CAM", 1),
        "Synergistic with everyone:",
        ("Sebas: tells him if VG2 sweep range is sufficient → device specs", 1),
        ("Robert: his differentiable emulator is what makes meta-plasticity learnable", 1),
        ("Mario: a clean story he can tell ('a chip that organizes itself')", 1),
        "Addresses the variability angle:",
        ("5–10% device spread is fine when the network can REROUTE around bad cells", 1),
        "Classical formulation: meta-plasticity (plasticity of plasticity).",
        ("Well-established in neuroscience. New for hardware.", 1),
    ])

    add_bullets_slide(prs, "Slide 27 — How Far Are We With It?", [
        "Honestly: idea stage.",
        "We have:",
        ("Working differentiable cell_fast", 1),
        ("Six topology generators", 1),
        ("Multi-task setup ready (Hopfield + reservoir together)", 1),
        ("Just verified V_latch is INERT in our model", 1),
        "Next sanity check (5 min compute):",
        ("Sweep one cell's VG2 from −0.4 to +0.5", 1),
        ("Plot Vb(t) under a drive pulse, confirm three regimes from Slide 11", 1),
        "If yes → run differentiable meta-plasticity training",
        "If no → cell_fast broken for this purpose, need Robert's emulator first",
    ])

    # ── PART VII ──
    add_section_slide(prs, "VII", "Asks and Proposed Deliverables")

    add_bullets_slide(prs, "Slide 28 — Three Things We Want to Discuss", [
        "For Sebastian:",
        ("Is the 130nm PTM card sufficient, or are foundry-specific tweaks needed?", 1),
        ("VG2 polarity: our model says +VG2 makes latching easier; data says opposite", 1),
        ("Could we get transient measurements (Id vs t under write pulse) for τ-calibration?", 1),
        "For Robert:",
        ("Could we use your differentiable cell emulator for network experiments? Co-authorship offered.", 1),
        ("Weekly cadence calls?", 1),
        "For Mario:",
        ("Vendor registration status — anything we can do to unblock?", 1),
        ("Does the meta-plasticity story fit the strategic narrative?", 1),
    ])

    add_bullets_slide(prs, "Slide 29 — Three Things We Propose to Deliver", [
        "1. Standardized network-level benchmark suite (low novelty, high utility)",
        ("Hopfield capacity + noise, Memory Capacity, NARMA, classification", 1),
        ("Run on Robert's emulator → reference numbers Seb can compare future cells against", 1),
        "2. Variability-aware design rules (uses NS-RAM's real strength)",
        ("Sweep ±5–10% per-cell parameter spread", 1),
        ("Output: 'for task X, you tolerate Y% spread before performance collapses'", 1),
        "3. Self-organizing meta-plasticity (high novelty)",
        ("Per-cell learnable VG2 trained on multi-task signal", 1),
        ("If it works, it's a paper-grade demonstration", 1),
        "Start with #1 (have infrastructure), use #2 for design recs, #3 as speculative-but-exciting.",
    ])

    add_diagram_slide(prs, "Slide 30 — Roles in One Page",
"""    SEBASTIAN  — makes and measures the cell.
                  Owns: physical wafers, I-V data, foundry SPICE deck.

    ROBERT     — fast, differentiable copy of the cell.
                  Owns: NN emulator (~30,000× faster than SPICE).

    MARIO      — strategy, funding, narrative.
                  Owns: PI role, publications, foundry relationships.

    ERIC + US  — what computes when many cells are wired together.
                  Owns: nsram package, network experiments, plasticity rules,
                        differentiable canonical model, meta-plasticity demo.

    ┌───────────────────────────────────────────┐
    │  NETWORK / APPLICATION   ← us             │
    ├───────────────────────────────────────────┤
    │  FAST CELL EMULATOR      ← Robert         │
    ├───────────────────────────────────────────┤
    │  PHYSICAL CELL + SPICE   ← Sebastian      │
    └───────────────────────────────────────────┘
        (Mario sets strategy around it all)
""",
        footer="Complementary, not competing. Each layer needs the others to deliver value.")

    add_bullets_slide(prs, "Slide 31 — Our Ask, Simply", [
        "To SEBASTIAN:",
        ("Confirm VG2 polarity (our model says +VG2 latches easier; your data says opposite)", 1),
        ("Send transient measurements (Id vs t under write pulse) for τ-calibration", 1),
        ("Weekly cadence calls so we stay aligned with new measurements", 1),
        "To ROBERT:",
        ("Could we plug your cell emulator into our network experiments? Co-authorship offered", 1),
        ("Would you co-validate our diff_canonical against your SPICE-trained NN?", 1),
        "To MARIO:",
        ("Vendor registration status — anything we can do to unblock?", 1),
        ("Does the meta-plasticity story (self-organized VG2 allocation) fit the strategic narrative?", 1),
        "Bottom line: we want to be useful network-level. We need your cell-level inputs to be credible.",
    ])

    add_bullets_slide(prs, "Slide 32 — What Success Looks Like + Closing", [
        "Six months from now, we'd ideally have:",
        ("✓ Reproducible benchmark suite published with results on Robert's emulator", 1),
        ("✓ Short technical note: 'Cell parameter ranking for NS-RAM array tasks'", 1),
        ("✓ One co-authored paper draft on meta-plasticity (if it works)", 1),
        ("✓ Updated nsram package consumed by Sebas's group as a standard tool", 1),
        ("✓ Clear input to Sebas for next fab iteration", 1),
        "What we are NOT trying to do:",
        ("Design or tape out chips ourselves (Mario's space)", 1),
        ("Build a competing cell model (Robert's space)", 1),
        ("Replace SPICE for cell-level work (Sebas's space)", 1),
        "Looking forward to the call on the 30th. — Eric & Claude",
    ])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}  ({OUT.stat().st_size / 1024:.0f} kB, "
          f"{len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
